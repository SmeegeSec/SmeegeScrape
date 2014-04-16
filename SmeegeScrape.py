#!/usr/bin/env python
"""
Name:           SmeegeScrape
Version:        0.6
Date:           01/27/2014
Author:         Smeege
Contact:        SmeegeSec@gmail.com

Description:    SmeegeScrape.py is a simple python script to scrape text from various sources including local files and 
                web pages, and turn the text into a custom word list.  A customized word list has many uses, from web 
                application testing to password cracking, having a specific set of words to use against a target can 
                increase efficiency and effectiveness during a penetration test.  I realize there are other text scrapers 
                publicly available however I feel this script is simple, efficient, and specific enough to warrant its own 
                release. This script is able to read almost any file which has cleartext in it that python can open.  I 
                have also included support for file formats such as pdf, html, docx, and pptx. 

Copyright (c) 2014, Smeege Sec (http://www.smeegesec.com)
All rights reserved.
Please see the attached LICENSE file for additional licensing information.
"""

import argparse
import docx #https://pypi.python.org/pypi/python-docx
import pptx #https://pypi.python.org/pypi/python-pptx
import glob
import mimetypes
import nltk #http://pypi.python.org/pypi/nltk
import os
import PyPDF2 #https://pypi.python.org/pypi/PyPDF2/1.19
import re
import string
from argparse import RawTextHelpFormatter
from collections import OrderedDict
from urllib import quote
from urllib2 import urlopen

title = "\n===SmeegeSec's========================================= \n"
title += "   ____                       ____                      \n"
title += "  / __/_ _  ___ ___ ___ ____ / __/__________ ____  ___  \n"
title += " _\ \/  ' \/ -_) -_) _ `/ -_)\ \/ __/ __/ _ `/ _ \/ -_) \n"
title += "/___/_/_/_/\__/\__/\_, /\__/___/\__/_/  \_,_/ .__/\__/  \n"
title += "                  /___/                    /_/          \n"
title += "================================================v0.6=== \n"
title += "\n"
print title

parser = argparse.ArgumentParser(description=title,formatter_class=RawTextHelpFormatter)

operationGroup = parser.add_mutually_exclusive_group(required=True)
operationGroup.add_argument('-f',  action="store", dest="localFile", help="Specify a local file to scrape.")
operationGroup.add_argument('-d', action="store", dest="fileDirectory", help="Specify a directory to scrape all supported files.")
operationGroup.add_argument('-u',  action="store", dest="webUrl", help="Specify a url to scrape page content (correct format: http(s)://www.smeegesec.com)")
operationGroup.add_argument('-l', action="store", dest="webList", help="Specify a text file with a list of URLs to scrape (separated by newline).")

optionGroup = parser.add_argument_group('paramters and options')
optionGroup.add_argument('-o', action="store", dest="outputFile", help="Output filename. (Default: smeegescrape_out.txt)")
optionGroup.add_argument('-i', action="store_true", dest="integers", help="Remove integers [0-9] from the final output.")
optionGroup.add_argument('-s', action="store_true", dest="specials", help="Remove special characters (only alphanum allowed) from the final output.")
optionGroup.add_argument('-n', action="store_true", dest="nonprintable", help="Remove all non-printable ASCII characters from the final output.")
optionGroup.add_argument('-min', action="store", dest="minLength", type=int, help="Set the minimum number of characters for each word (Default: 3).")
optionGroup.add_argument('-max', action="store", dest="maxLength", type=int, help="Set the maximum number of characters for each word (Default: 30).")

args = parser.parse_args()

#PyPDF2 does not always work very well with extracting text from pdf.  It mainly depends on how the PDF was generated.                
def getPDFContent(path):
    # Load PDF into PyPDF2
    pdf = PyPDF2.PdfFileReader(file(path, "rb"))

    if pdf.isEncrypted:
        print 'pdf ' + path + ' is encrypted, trying blank password...'
        pdf.decrypt('') #If you want to provide your own password for an encrypted pdf, modify code here.
    
    content = ""
    # Iterate PDF pages
    for i in range(0, pdf.getNumPages()):
        # Extract text from page and add to content
        content += pdf.getPage(i).extractText() + "\n"
    # Collapse whitespace
    content = " ".join(content.replace(u"\xa0", " ").strip().split())
    content = content.encode("ascii", "ignore")
    content = content.split(' ')
    if args.minLength or args.maxLength:
        for word in content:
            if not(len(word.translate(None,charBlacklist)) < minl or len(word) > maxl):
                wordList.append(str(word).translate(None,charBlacklist))
    else:
        for word in content:
            wordList.append(str(word).translate(None,charBlacklist))

def webUrl(fullUrl):
    #urllib2 works best with a specific url format
    validUrl = re.compile(
        r'^(?:http)s?://|' # http:// or https://
        r'^(?:http)s?://www.'
        r'(?:(?:[A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+(?:[A-Z]{2,6}\.?|[A-Z0-9-]{2,}\.?)|' #domain...
        r'localhost|' #localhost...
        r'\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})' # or ip
        r'(?::\d+)?' # optional port
        r'(?:/?|[/?]\S+)$', re.IGNORECASE)

    if validUrl.match(fullUrl):
        finalList = []
        urlInput = quote(fullUrl, safe="%/:=&?~#+!$,;'@()*[]")
        urlInput = urlInput.strip('%0A')
        try:
            u = urlopen(urlInput)
            html = u.read()
            raw = nltk.clean_html(html)
            tokens = nltk.word_tokenize(raw)
            if args.minLength or args.maxLength:
                for token in tokens:
                    if not(len(token.translate(None,charBlacklist)) < minl or len(token) > maxl):
                        wordList.append(str(token).translate(None,charBlacklist))
            else:
                for token in tokens:
                    wordList.append(str(token).translate(None,charBlacklist))
            print "Scraping URL - {0}".format(fullUrl)
        except Exception as e:
            print 'There was an error connecting to or parsing {0}'.format(fullUrl)
            print 'Error: %s' % e
    else:
        print 'INVALID URL - {0}. Format must be http(s)://www.smeegesec.com.'.format(fullUrl)

def webList(webListFile):
    if os.path.isfile(webListFile):
        with open(webListFile) as f:
            webList = f.readlines()

        for url in webList:
            webUrl(url.rstrip('\n'))

        f.close()
    else:
        print 'Error opening file: {0}'.format(fileInput)

def localFile(fileInput):
    if os.path.isfile(fileInput):
        print "Scraping Local File - {0}".format(fileInput)
        mimetypes.init()
        file_type, file_encoding = mimetypes.guess_type(fileInput)
        print file_type
        if file_type == 'application/pdf':
            getPDFContent(fileInput)
        elif file_type == 'text/html':
            raw = nltk.clean_html(open(fileInput).read())
            tokens = nltk.word_tokenize(raw)
            if args.minLength or args.maxLength:
                for token in tokens:
                    if not(len(token.translate(None,charBlacklist)) < minl or len(token) > maxl):
                        wordList.append(str(token).translate(None,charBlacklist))
            else:
                for token in tokens:
                    wordList.append(str(token).translate(None,charBlacklist))
        elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            document = docx.opendocx(fileInput)
            sentances = docx.getdocumenttext(document)
            sentances = map(lambda s: s.encode("ascii", "ignore"), sentances)
            if args.minLength or args.maxLength:
                for sentance in sentances:
                    for word in set(sentance.split()):
                        if not(len(str(word).translate(None,charBlacklist)) < minl or len(str(word)) > maxl):
                            wordList.append(str(word).translate(None,charBlacklist))
            else:
                for sentance in sentances:
                    for word in set(sentance.split()):
                        wordList.append(str(word).translate(None,charBlacklist))
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or file_type == 'application/x-mspowerpoint.12':
            try:
                prs = pptx.Presentation(fileInput)
                text_runs = list()

                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_textframe:
                            continue
                        for paragraph in shape.textframe.paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text)

                if args.minLength or args.maxLength:
                    for sentance in text_runs:
                        for word in set(sentance.split()):
                            if not(len(str((word.translate(None,charBlacklist)))) < minl or len(str(word)) > maxl):
                                wordList.append(str(word).translate(None,charBlacklist))
                else:
                    for sentance in text_runs:
                        for word in set(sentance.split()):
                            wordList.append(str(word).translate(None,charBlacklist))
            except Exception as e:
                print 'Error opening file: {0}'.format(fileInput)
                pass
        else: #'text/plain' or unknown format
            try:
                words = set(open(fileInput).read().split())

                if args.minLength or args.maxLength:
                    for word in words:
                        if not(len(str((word.translate(None,charBlacklist)))) < minl or len(str(word)) > maxl):
                            wordList.append(str(word).translate(None,charBlacklist))                  
                else:
                    for word in words:
                        wordList.append(str(word).translate(None,charBlacklist))
            except:
                print 'Error opening file: {0}'.format(fileInput)
                pass
    else:
        print 'Error opening file: {0}'.format(fileInput)

def fileDir(dirInput):
    if os.path.isdir(dirInput):
        "Scraping Files in Local Directory - {0}".format(dirInput)
        mimetypes.init()
        fileTypes = ('*.pdf', '*.docx', '*.txt', '*.csv', '*.py', '*.cpp', '*.pl', '*.log', '*.rtf', '*.css', '*.dat', '*.html', '*.php', '*.pps', '*.ppt', '*.pptx', '*.sh', '*.xml', '*.xsl')
        for fileType in fileTypes:
            for globFile in glob.glob(os.path.join(dirInput, fileType)):
                file_type, file_encoding = mimetypes.guess_type(globFile)
                localFile(globFile)
    else:
        print 'Error accessing directory: {0}'.format(dirInput)

def output():
    try:
        if not args.outputFile:
            args.outputFile = 'smeegescrape_out.txt'
        outputFile = open(args.outputFile, 'w')
        wordListFinal = OrderedDict.fromkeys(wordList).keys()
        
        for word in wordListFinal:
            outputFile.write(word)
            outputFile.write('\n')
        outputFile.close()
        
        print '\n{0} unique words have been scraped.'.format(len(wordListFinal))
        print 'Output file successfully written: {0}'.format(outputFile.name)
    except Exception as e:
        print 'Error creating output file: {0}'.format(outputFile.name)
        print e

if __name__ == "__main__":

    wordList = list()
    charBlacklist = ""

    if args.minLength or args.maxLength:
        minl = args.minLength if args.minLength else 3
        maxl = args.maxLength if args.maxLength else 30
        if minl > maxl:
            print 'Argument minLength cannot be greater than maxLength. Setting defaults to min=3 max=30.'
            minl = 3
            maxl = 30

    charBlacklist = ""

    if args.specials:
        charBlacklist += "~`!@#$%^&*()_\"--+=[]{}|/\:;'<,>.?/"
    if args.integers:
        charBlacklist += string.digits
    if args.nonprintable:
        non_printable = filter(lambda x: x not in string.printable, map(chr, range(0, 256)))
        charBlacklist += ''.join(non_printable)

    if args.localFile:
        localFile(args.localFile)
    if args.fileDirectory:
        fileDir(args.fileDirectory)
    if args.webUrl:
        webUrl(args.webUrl)
    if args.webList:
        webList(args.webList)

    output()
    
